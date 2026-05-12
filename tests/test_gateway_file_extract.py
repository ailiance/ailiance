"""Tests for the /v1/files/extract endpoint + extract module."""

from __future__ import annotations

import io
from typing import Any

import pytest
from fastapi.testclient import TestClient

from src.gateway.file_extract import (
    ExtractError,
    detect_format,
    extract,
)


# ---------------------------------------------------------------------------
# Helpers — generate the smallest valid file we can for each format.
# ---------------------------------------------------------------------------


def _make_docx(text: str = "Hello world") -> bytes:
    from docx import Document

    buf = io.BytesIO()
    doc = Document()
    doc.add_heading("Title", level=1)
    doc.add_paragraph(text)
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx(rows: list[list[Any]]) -> bytes:
    from openpyxl import Workbook

    buf = io.BytesIO()
    wb = Workbook()
    ws = wb.active
    ws.title = "S1"
    for row in rows:
        ws.append(row)
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(titles: list[str]) -> bytes:
    from pptx import Presentation

    buf = io.BytesIO()
    prs = Presentation()
    for t in titles:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = t
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# detect_format
# ---------------------------------------------------------------------------


class TestDetectFormat:
    def test_mime_pdf(self):
        assert detect_format(None, "application/pdf") == "pdf"

    def test_mime_with_charset(self):
        # Browsers append ;charset=utf-8 — we must strip it.
        assert detect_format(None, "text/plain; charset=utf-8") == "text"

    def test_filename_fallback(self):
        # MIME unhelpful → ext drives detection.
        assert detect_format("report.PDF", "application/octet-stream") == "pdf"

    def test_unsupported_returns_none(self):
        assert detect_format("a.exe", "application/octet-stream") is None

    def test_no_hints_returns_none(self):
        assert detect_format(None, None) is None


# ---------------------------------------------------------------------------
# extract — per-format
# ---------------------------------------------------------------------------


class TestExtract:
    def test_text_passthrough(self):
        r = extract(b"hello\nworld", filename="a.txt", mime="text/plain")
        assert r.format == "text"
        assert "hello" in r.markdown and "world" in r.markdown

    def test_html_strips_scripts(self):
        html = b"<html><head><title>T</title></head><body><script>alert(1)</script><p>visible</p></body></html>"
        r = extract(html, filename="a.html", mime="text/html")
        assert r.format == "html"
        assert "visible" in r.markdown
        assert "alert" not in r.markdown
        assert r.metadata.get("title") == "T"

    def test_docx_emits_heading_and_body(self):
        data = _make_docx("body text here")
        r = extract(data, filename="a.docx")
        assert r.format == "docx"
        assert "Title" in r.markdown
        assert "body text here" in r.markdown
        # Heading 1 → '# '
        assert r.markdown.startswith("# ")

    def test_xlsx_emits_sheet_header_and_rows(self):
        data = _make_xlsx([["a", "b"], [1, 2]])
        r = extract(data, filename="a.xlsx")
        assert r.format == "xlsx"
        assert "Sheet: S1" in r.markdown
        assert "1 | 2" in r.markdown

    def test_pptx_emits_slide_headers(self):
        data = _make_pptx(["First", "Second"])
        r = extract(data, filename="a.pptx")
        assert r.format == "pptx"
        assert "Slide 1" in r.markdown
        assert "First" in r.markdown
        assert "Second" in r.markdown

    def test_unsupported_raises(self):
        with pytest.raises(ExtractError) as exc_info:
            extract(b"not a real file", filename="a.bin", mime="application/octet-stream")
        assert exc_info.value.code == "unsupported_format"

    def test_corrupt_pdf_raises_extract_failed(self):
        with pytest.raises(ExtractError) as exc_info:
            extract(b"not a pdf", filename="x.pdf", mime="application/pdf")
        assert exc_info.value.code == "extract_failed"

    def test_oversize_raises_file_too_large(self):
        # 26 MiB of zeros — over the 25 MiB cap.
        with pytest.raises(ExtractError) as exc_info:
            extract(b"\x00" * (26 * 1024 * 1024), filename="big.txt", mime="text/plain")
        assert exc_info.value.code == "file_too_large"

    def test_truncation_metadata_when_huge_markdown(self):
        # 250k chars text → truncated to MAX_MARKDOWN_CHARS (200k) with
        # metadata.truncated set.
        big_text = ("a" * 1000 + "\n") * 250  # ~250k chars
        r = extract(big_text.encode(), filename="big.txt", mime="text/plain")
        assert r.metadata.get("truncated") is True
        assert r.metadata.get("original_chars", 0) > r.metadata.get("kept_chars", 0)
        assert len(r.markdown) == r.metadata["kept_chars"]


# ---------------------------------------------------------------------------
# HTTP endpoint
# ---------------------------------------------------------------------------


@pytest.fixture
def client():
    from src.gateway.server import make_gateway_app

    return TestClient(make_gateway_app(skip_router_load=True))


class TestExtractEndpoint:
    def test_missing_file_returns_400(self, client):
        resp = client.post("/v1/files/extract")
        assert resp.status_code == 400
        assert resp.json()["detail"]["code"] == "missing_file"

    def test_unsupported_mime_returns_400(self, client):
        resp = client.post(
            "/v1/files/extract",
            files={"file": ("a.bin", b"x", "application/octet-stream")},
        )
        assert resp.status_code == 400
        assert resp.json()["detail"]["code"] == "unsupported_format"

    def test_docx_roundtrip(self, client):
        data = _make_docx("hello from docx")
        resp = client.post(
            "/v1/files/extract",
            files={
                "file": (
                    "report.docx",
                    data,
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            },
        )
        assert resp.status_code == 200, resp.text
        body = resp.json()
        assert body["format"] == "docx"
        assert "hello from docx" in body["markdown"]
        assert body["filename"] == "report.docx"

    def test_xlsx_roundtrip(self, client):
        data = _make_xlsx([["alpha", "beta"], [42, 100]])
        resp = client.post(
            "/v1/files/extract",
            files={"file": ("data.xlsx", data, "")},  # empty mime → ext fallback
        )
        assert resp.status_code == 200, resp.text
        body = resp.json()
        assert body["format"] == "xlsx"
        assert "alpha" in body["markdown"]
        assert "42 | 100" in body["markdown"]
