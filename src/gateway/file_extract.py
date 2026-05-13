"""File → markdown extraction for ``/v1/files/extract``.

Lightweight, deps-light alternative to markitdown's full stack (we
explicitly avoid the onnxruntime/whisper extras that markitdown's
opinionated extras pull). Each format gets a small dedicated extractor;
the dispatcher routes by MIME type, then by filename extension as a
fallback when MIME is ``application/octet-stream`` (common for browser
uploads of office files).

All extractors return ``(markdown_text, metadata)`` and never raise to
the caller; failures bubble up as :class:`ExtractError` with a structured
``code`` so the endpoint can translate to a clean HTTP 4xx/5xx.
"""

from __future__ import annotations

import io
import logging
import re
from dataclasses import dataclass, field
from typing import Any

log = logging.getLogger(__name__)


class ExtractError(Exception):
    """Raised when a file cannot be extracted to text."""

    def __init__(self, code: str, message: str) -> None:
        super().__init__(message)
        self.code = code
        self.message = message


# Hard limit on uploaded bytes — large files are usually accidents
# (DOC archives, video, etc.) and the goal here is to feed an LLM, not
# to be a general document store.
MAX_BYTES = 25 * 1024 * 1024  # 25 MiB

# Maximum markdown characters returned — caps prompt cost downstream.
# Truncation is signalled in metadata so callers can decide how to react.
MAX_MARKDOWN_CHARS = 200_000


@dataclass
class ExtractResult:
    markdown: str
    format: str  # canonical short tag: "pdf" | "docx" | "xlsx" | "pptx" | "text" | "html"
    metadata: dict[str, Any] = field(default_factory=dict)


_EXT_TO_FORMAT = {
    "pdf": "pdf",
    "docx": "docx",
    "xlsx": "xlsx",
    "pptx": "pptx",
    "txt": "text",
    "md": "text",
    "markdown": "text",
    "html": "html",
    "htm": "html",
    # Image OCR
    "png": "image",
    "jpg": "image",
    "jpeg": "image",
    "gif": "image",
    "bmp": "image",
    "webp": "image",
    "tiff": "image",
    "tif": "image",
}

_MIME_TO_FORMAT = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "text/plain": "text",
    "text/markdown": "text",
    "text/html": "html",
    "image/png": "image",
    "image/jpeg": "image",
    "image/jpg": "image",
    "image/gif": "image",
    "image/bmp": "image",
    "image/webp": "image",
    "image/tiff": "image",
}

# Tesseract language codes — set to French + English so the OCR
# pipeline handles both without per-call language hints. Add more
# (e.g. "spa", "deu") by overriding via OCR_LANGS env var.
_OCR_LANGS = "fra+eng"

# Minimum number of characters pdfminer must extract before we trust
# the PDF as text-bearing. Below this we fall back to OCR (page-by-
# page rasterisation + Tesseract). 50 chars is roughly one line of
# body text — empty / scanned / image-only PDFs typically yield 0-10.
_PDF_OCR_THRESHOLD = 50


def detect_format(filename: str | None, mime: str | None) -> str | None:
    """Return canonical short format tag or ``None`` if unsupported."""
    if mime:
        fmt = _MIME_TO_FORMAT.get(mime.split(";")[0].strip().lower())
        if fmt:
            return fmt
    if filename:
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        return _EXT_TO_FORMAT.get(ext)
    return None


def extract(data: bytes, *, filename: str | None = None, mime: str | None = None) -> ExtractResult:
    """Extract markdown text from ``data``.

    Raises :class:`ExtractError` on unsupported format or extractor failure.
    Truncates the markdown body to :data:`MAX_MARKDOWN_CHARS` and records
    the original length in metadata when truncated.
    """
    if len(data) > MAX_BYTES:
        raise ExtractError(
            "file_too_large",
            f"File exceeds {MAX_BYTES // (1024 * 1024)} MiB limit "
            f"(got {len(data) // (1024 * 1024)} MiB).",
        )
    fmt = detect_format(filename, mime)
    if not fmt:
        raise ExtractError(
            "unsupported_format",
            f"Could not detect a supported format from filename={filename!r} mime={mime!r}. "
            "Supported: PDF, DOCX, XLSX, PPTX, TXT, MD, HTML, PNG/JPG/GIF/BMP/WEBP/TIFF.",
        )

    extractor = _DISPATCH[fmt]
    try:
        md, meta = extractor(data)
    except ExtractError:
        raise
    except Exception as exc:  # noqa: BLE001
        log.warning("extract %s failed: %s", fmt, exc, exc_info=True)
        raise ExtractError(
            "extract_failed",
            f"{fmt} extractor raised: {type(exc).__name__}: {exc}",
        ) from exc

    full_len = len(md)
    if full_len > MAX_MARKDOWN_CHARS:
        md = md[:MAX_MARKDOWN_CHARS]
        meta["truncated"] = True
        meta["original_chars"] = full_len
        meta["kept_chars"] = MAX_MARKDOWN_CHARS
    return ExtractResult(markdown=md, format=fmt, metadata=meta)


# ---------------------------------------------------------------------------
# Per-format extractors. Each returns (markdown, metadata_dict).
# ---------------------------------------------------------------------------

_WS_RUN = re.compile(r"[ \t]+\n")


def _normalize(text: str) -> str:
    """Collapse trailing whitespace and triple-newlines."""
    text = _WS_RUN.sub("\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _extract_text(data: bytes) -> tuple[str, dict[str, Any]]:
    return data.decode("utf-8", errors="replace"), {}


def _extract_html(data: bytes) -> tuple[str, dict[str, Any]]:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(data, "html.parser")
    # Drop script/style noise before extracting visible text.
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    return _normalize(text), {"title": (soup.title.string if soup.title else None)}


def _extract_pdf(data: bytes) -> tuple[str, dict[str, Any]]:
    """Extract PDF text.

    Tries pdfminer first (fast, no extra binaries). If the resulting
    text is too short to be a real document body — typical for
    scanned / image-only PDFs — falls back to OCR: rasterise each
    page via :mod:`pdf2image` (poppler) then Tesseract over every
    page image. ``metadata.ocr=True`` flags the OCR path so callers
    can warn about lower confidence.
    """
    from pdfminer.high_level import extract_text

    text = extract_text(io.BytesIO(data))
    text = _normalize(text)
    pages = text.count("\f") + 1 if text else 0
    if len(text) >= _PDF_OCR_THRESHOLD:
        return text, {"pages": pages, "ocr": False}

    # OCR fallback for scanned PDFs.
    try:
        ocr_text, ocr_pages = _ocr_pdf(data)
    except ExtractError:
        raise
    except Exception as exc:  # noqa: BLE001
        log.warning("PDF OCR fallback failed: %s", exc, exc_info=True)
        # Return whatever pdfminer gave us rather than erroring — the
        # caller will at least see "no text found" instead of HTTP 4xx
        # for a slightly malformed PDF.
        return text, {"pages": pages, "ocr": False, "ocr_attempted": True}
    return ocr_text, {"pages": ocr_pages, "ocr": True}


def _ocr_pdf(data: bytes) -> tuple[str, int]:
    """Rasterise + OCR each page. Returns ``(text, page_count)``."""
    from pdf2image import convert_from_bytes
    import pytesseract

    images = convert_from_bytes(data, dpi=200)
    parts: list[str] = []
    for i, img in enumerate(images, start=1):
        page_text = pytesseract.image_to_string(img, lang=_OCR_LANGS)
        parts.append(f"## Page {i}")
        parts.append(_normalize(page_text))
    return _normalize("\n\n".join(parts)), len(images)


def _extract_image(data: bytes) -> tuple[str, dict[str, Any]]:
    """OCR a standalone image upload (PNG / JPEG / etc.).

    Returns the recognised text as markdown. Empty result is a valid
    outcome — Tesseract often returns no characters on photo
    backgrounds, drawings, or pure-graphic images.
    """
    import pytesseract
    from PIL import Image

    img = Image.open(io.BytesIO(data))
    # Convert palette / mode quirks to RGB so Tesseract has a stable
    # input. PNG with alpha would otherwise raise.
    if img.mode not in ("RGB", "L"):
        img = img.convert("RGB")
    text = pytesseract.image_to_string(img, lang=_OCR_LANGS)
    return _normalize(text), {
        "width": img.width,
        "height": img.height,
        "format": img.format or "unknown",
        "ocr": True,
    }


def _extract_docx(data: bytes) -> tuple[str, dict[str, Any]]:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    for para in doc.paragraphs:
        style = (para.style.name if para.style else "") or ""
        text = para.text.strip()
        if not text:
            continue
        if style.startswith("Heading"):
            try:
                level = int(style.replace("Heading", "").strip() or "1")
            except ValueError:
                level = 1
            parts.append(f"{'#' * min(level, 6)} {text}")
        else:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return _normalize("\n\n".join(parts)), {"paragraphs": len(doc.paragraphs)}


def _extract_xlsx(data: bytes) -> tuple[str, dict[str, Any]]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    parts: list[str] = []
    sheet_count = 0
    for sheet in wb.worksheets:
        sheet_count += 1
        parts.append(f"## Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c for c in cells):
                parts.append(" | ".join(cells))
        parts.append("")
    return _normalize("\n".join(parts)), {"sheets": sheet_count}


def _extract_pptx(data: bytes) -> tuple[str, dict[str, Any]]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
        parts.append("")
    return _normalize("\n".join(parts)), {"slides": len(prs.slides)}


_DISPATCH = {
    "pdf": _extract_pdf,
    "docx": _extract_docx,
    "xlsx": _extract_xlsx,
    "pptx": _extract_pptx,
    "text": _extract_text,
    "html": _extract_html,
    "image": _extract_image,
}
